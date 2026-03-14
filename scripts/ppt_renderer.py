
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE

def create_ppt():
    # 1. Define JSON Data based on Visual Analysis
    data = {
        "slide": {
            "width_inches": 13.33,
            "height_inches": 7.5,
            "background_color": "D9EDF7" # Light Blue background
        },
        "elements": [
            # --- Top Left: Users ---
            {
                "type": "SHAPE", "shape_type": "ROUNDED_RECTANGLE", "text": "",
                "position": {"x": 0.5, "y": 0.5, "w": 1.5, "h": 3.2},
                "style": {"fill_color": None, "border_color": "000000", "is_dashed": True}
            },
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "业务\n部门",
                "position": {"x": 0.7, "y": 0.7, "w": 1.1, "h": 0.5},
                "style": {"fill_color": "E2F0D9", "font_size": 10, "border_color": "5E9CD3"}
            },
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "委办\n局",
                "position": {"x": 0.7, "y": 1.3, "w": 1.1, "h": 0.5},
                "style": {"fill_color": "E2F0D9", "font_size": 10, "border_color": "5E9CD3"}
            },
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "用人\n单位",
                "position": {"x": 0.7, "y": 1.9, "w": 1.1, "h": 0.5},
                "style": {"fill_color": "E2F0D9", "font_size": 10, "border_color": "5E9CD3"}
            },
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "广大\n公众",
                "position": {"x": 0.7, "y": 2.5, "w": 1.1, "h": 0.5},
                "style": {"fill_color": "E2F0D9", "font_size": 10, "border_color": "5E9CD3"}
            },

            # --- Arrow: Users <- Portals ---
            {
                "type": "SHAPE", "shape_type": "LEFT_ARROW", "text": "",
                "position": {"x": 2.2, "y": 1.5, "w": 1.0, "h": 1.2},
                "style": {"fill_color": "4472C4", "border_color": "2F5597"}
            },

            # --- Top Middle-Left: Portals ---
            {
                "type": "SHAPE", "shape_type": "ROUNDED_RECTANGLE", "text": "",
                "position": {"x": 3.4, "y": 0.5, "w": 1.5, "h": 3.2},
                "style": {"fill_color": None, "border_color": "000000", "is_dashed": True}
            },
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "内网\n门户",
                "position": {"x": 3.6, "y": 0.9, "w": 1.1, "h": 1.0},
                "style": {"fill_color": "E4D7F5", "font_size": 12, "border_color": "7030A0"}
            },
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "外网\n门户",
                "position": {"x": 3.6, "y": 2.1, "w": 1.1, "h": 1.0},
                "style": {"fill_color": "E4D7F5", "font_size": 12, "border_color": "7030A0"}
            },

            # --- Arrow: Portals <- Analysis ---
            {
                "type": "SHAPE", "shape_type": "LEFT_ARROW", "text": "",
                "position": {"x": 5.1, "y": 1.5, "w": 1.0, "h": 1.2},
                "style": {"fill_color": "4472C4", "border_color": "2F5597"}
            },

            # --- Top Middle-Right: Analysis ---
            {
                "type": "SHAPE", "shape_type": "ROUNDED_RECTANGLE", "text": "数据分析与展现",
                "position": {"x": 6.3, "y": 0.5, "w": 4.5, "h": 3.2},
                "style": {"fill_color": "E1F0FA", "border_color": "000000", "is_dashed": True, "alignment": "TOP", "font_size": 14}
            },
            # Bars (7 items)
            # "数据查询", "统计管理", "资源汇总", "统一分析", "报表管理", "预测预警", "决策支持"
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "数据\n查询",
                "position": {"x": 6.5, "y": 1.2, "w": 0.4, "h": 2.2},
                "style": {"fill_color": "F8CBAD", "font_size": 11, "border_color": "C65911"}
            },
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "统计\n管理",
                "position": {"x": 7.05, "y": 1.2, "w": 0.4, "h": 2.2},
                "style": {"fill_color": "F8CBAD", "font_size": 11, "border_color": "C65911"}
            },
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "资源\n汇总",
                "position": {"x": 7.6, "y": 1.2, "w": 0.4, "h": 2.2},
                "style": {"fill_color": "F8CBAD", "font_size": 11, "border_color": "C65911"}
            },
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "统一\n分析",
                "position": {"x": 8.15, "y": 1.2, "w": 0.4, "h": 2.2},
                "style": {"fill_color": "F8CBAD", "font_size": 11, "border_color": "C65911"}
            },
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "报表\n管理",
                "position": {"x": 8.7, "y": 1.2, "w": 0.4, "h": 2.2},
                "style": {"fill_color": "F8CBAD", "font_size": 11, "border_color": "C65911"}
            },
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "预测\n预警",
                "position": {"x": 9.25, "y": 1.2, "w": 0.4, "h": 2.2},
                "style": {"fill_color": "F8CBAD", "font_size": 11, "border_color": "C65911"}
            },
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "决策\n支持",
                "position": {"x": 9.8, "y": 1.2, "w": 0.4, "h": 2.2},
                "style": {"fill_color": "F8CBAD", "font_size": 11, "border_color": "C65911"}
            },

            # --- Arrow: Analysis <- Platform ---
            {
                "type": "SHAPE", "shape_type": "LEFT_ARROW", "text": "",
                "position": {"x": 11.0, "y": 1.5, "w": 1.0, "h": 1.2},
                "style": {"fill_color": "4472C4", "border_color": "2F5597"}
            },

            # --- Top Right: Platform ---
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "数\n据\n整\n合\n与\n交\n换\n平\n台",
                "position": {"x": 12.2, "y": 0.5, "w": 0.9, "h": 3.2},
                "style": {"fill_color": "B4C7E7", "font_size": 14, "border_color": "5B9BD5"}
            },


            # --- Bottom Left: Apps ---
            {
                "type": "SHAPE", "shape_type": "ROUNDED_RECTANGLE", "text": "集成应用\n统一管理",
                "position": {"x": 0.5, "y": 4.2, "w": 2.0, "h": 3.0},
                "style": {"fill_color": None, "border_color": "000000", "is_dashed": True, "alignment": "TOP", "font_size": 12}
            },
            {
                "type": "SHAPE", "shape_type": "CHEVRON", "text": "原有应\n用系统",
                "position": {"x": 0.7, "y": 5.0, "w": 1.6, "h": 0.5},
                "style": {"fill_color": "5B9BD5", "font_size": 10, "border_color": "2F5597"}
            },
            {
                "type": "SHAPE", "shape_type": "CHEVRON", "text": "升级应\n用系统",
                "position": {"x": 0.7, "y": 5.7, "w": 1.6, "h": 0.5},
                "style": {"fill_color": "5B9BD5", "font_size": 10, "border_color": "2F5597"}
            },
            {
                "type": "SHAPE", "shape_type": "CHEVRON", "text": "新建应\n用系统",
                "position": {"x": 0.7, "y": 6.4, "w": 1.6, "h": 0.5},
                "style": {"fill_color": "5B9BD5", "font_size": 10, "border_color": "2F5597"}
            },

            # --- Arrows: Apps -> Resources ---
            {
                "type": "SHAPE", "shape_type": "RIGHT_ARROW", "text": "",
                "position": {"x": 2.7, "y": 5.1, "w": 0.6, "h": 0.3},
                "style": {"fill_color": "5B9BD5", "border_color": "2F5597"}
            },
            {
                "type": "SHAPE", "shape_type": "RIGHT_ARROW", "text": "",
                "position": {"x": 2.7, "y": 5.8, "w": 0.6, "h": 0.3},
                "style": {"fill_color": "5B9BD5", "border_color": "2F5597"}
            },
            {
                "type": "SHAPE", "shape_type": "RIGHT_ARROW", "text": "",
                "position": {"x": 2.7, "y": 6.5, "w": 0.6, "h": 0.3},
                "style": {"fill_color": "5B9BD5", "border_color": "2F5597"}
            },

            # --- Bottom Middle: Resources ---
            {
                "type": "SHAPE", "shape_type": "ROUNDED_RECTANGLE", "text": "",
                "position": {"x": 3.5, "y": 4.2, "w": 2.0, "h": 3.0},
                "style": {"fill_color": None, "border_color": "000000", "is_dashed": True}
            },
            {
                "type": "SHAPE", "shape_type": "CAN", "text": "结构化资\n源",
                "position": {"x": 3.8, "y": 4.5, "w": 1.4, "h": 1.0},
                "style": {"fill_color": "C5E0B4", "font_size": 11, "border_color": "548235"}
            },
            {
                "type": "SHAPE", "shape_type": "FOLDED_CORNER", "text": "非结构化\n资源",
                "position": {"x": 3.9, "y": 6.0, "w": 1.2, "h": 1.0},
                "style": {"fill_color": "FFF2CC", "font_size": 11, "border_color": "BF8F00"}
            },

            # --- Arrows: Resources -> Collection ---
            {
                "type": "SHAPE", "shape_type": "RIGHT_ARROW", "text": "",
                "position": {"x": 5.7, "y": 4.8, "w": 0.6, "h": 0.3},
                "style": {"fill_color": "5B9BD5", "border_color": "2F5597"}
            },
            {
                "type": "SHAPE", "shape_type": "RIGHT_ARROW", "text": "",
                "position": {"x": 5.7, "y": 6.4, "w": 0.6, "h": 0.3},
                "style": {"fill_color": "5B9BD5", "border_color": "2F5597"}
            },

            # --- Bottom Right: Collection ---
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "结构化数据采集",
                "position": {"x": 6.5, "y": 4.2, "w": 5.5, "h": 2.2},
                "style": {"fill_color": "BDD7EE", "border_color": "2F5597", "is_dashed": True, "alignment": "TOP", "font_size": 12}
            },
            # Inner Flow
            {
                "type": "SHAPE", "shape_type": "ROUNDED_RECTANGLE", "text": "采集接\n口管理",
                "position": {"x": 6.7, "y": 4.8, "w": 0.8, "h": 1.4},
                "style": {"fill_color": "E4D7F5", "font_size": 10, "border_color": "7030A0"}
            },
            {
                "type": "SHAPE", "shape_type": "RIGHT_ARROW", "text": "",
                "position": {"x": 7.6, "y": 5.3, "w": 0.4, "h": 0.4},
                "style": {"fill_color": "2F5597", "border_color": "2F5597"}
            },
            {
                "type": "SHAPE", "shape_type": "ROUNDED_RECTANGLE", "text": "采集模\n板定制",
                "position": {"x": 8.1, "y": 4.8, "w": 0.8, "h": 1.4},
                "style": {"fill_color": "E4D7F5", "font_size": 10, "border_color": "7030A0"}
            },
            {
                "type": "SHAPE", "shape_type": "RIGHT_ARROW", "text": "",
                "position": {"x": 9.0, "y": 5.3, "w": 0.4, "h": 0.4},
                "style": {"fill_color": "2F5597", "border_color": "2F5597"}
            },
            {
                "type": "SHAPE", "shape_type": "ROUNDED_RECTANGLE", "text": "采集资\n源审核",
                "position": {"x": 9.5, "y": 4.8, "w": 0.8, "h": 1.4},
                "style": {"fill_color": "E4D7F5", "font_size": 10, "border_color": "7030A0"}
            },
            {
                "type": "SHAPE", "shape_type": "RIGHT_ARROW", "text": "",
                "position": {"x": 10.4, "y": 5.3, "w": 0.4, "h": 0.4},
                "style": {"fill_color": "2F5597", "border_color": "2F5597"}
            },
            {
                "type": "SHAPE", "shape_type": "ROUNDED_RECTANGLE", "text": "数据处\n理分析",
                "position": {"x": 10.9, "y": 4.8, "w": 0.8, "h": 1.4},
                "style": {"fill_color": "E4D7F5", "font_size": 10, "border_color": "7030A0"}
            },
            
            # Bottom Bar
            {
                "type": "SHAPE", "shape_type": "RECTANGLE", "text": "非结构化数据采集工具",
                "position": {"x": 6.5, "y": 6.6, "w": 5.5, "h": 0.6},
                "style": {"fill_color": "9CC2E5", "font_size": 12, "border_color": "2F5597"}
            },

            # --- Arrow: Collection -> Platform ---
            {
                "type": "SHAPE", "shape_type": "UP_ARROW", "text": "",
                "position": {"x": 12.3, "y": 3.8, "w": 0.6, "h": 0.6},
                "style": {"fill_color": "4472C4", "border_color": "2F5597"}
            },
        ]
    }

    # 2. Render Logic
    prs = Presentation()
    prs.slide_width = Inches(data['slide']['width_inches'])
    prs.slide_height = Inches(data['slide']['height_inches'])
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank

    # Background Color
    bg_rgb = RGBColor.from_string(data['slide']['background_color'])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_rgb

    for el in data['elements']:
        if el['type'] == 'SHAPE':
            # Map shape type string to MSO_SHAPE enum
            # Default to RECTANGLE if not found
            shape_type_enum = getattr(MSO_SHAPE, el['shape_type'], MSO_SHAPE.RECTANGLE)
            
            x = Inches(el['position']['x'])
            y = Inches(el['position']['y'])
            w = Inches(el['position']['w'])
            h = Inches(el['position']['h'])
            
            shape = slide.shapes.add_shape(shape_type_enum, x, y, w, h)
            
            # Text
            if el.get('text'):
                shape.text = el['text']
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    text_frame.word_wrap = True
                    # Vertical alignment defaults to middle, check style
                    if el['style'].get('alignment') == 'TOP':
                         text_frame.vertical_anchor = MSO_ANCHOR.TOP
                    else:
                         text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                    for p in text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER
                        if el['style'].get('font_size'):
                            p.font.size = Pt(el['style']['font_size'])
                        p.font.name = 'Microsoft YaHei' # Font support
            
            # Style
            style = el.get('style', {})
            
            # Fill
            if style.get('fill_color'):
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(style['fill_color'])
            else:
                shape.fill.background() # No fill

            # Border
            line = shape.line
            if style.get('border_color'):
                line.color.rgb = RGBColor.from_string(style['border_color'])
                line.width = Pt(1.5)
            
            if style.get('is_dashed'):
                line.dash_style = MSO_LINE_DASH_STYLE.DASH
                
    
    output_file = 'data_integration_platform.pptx'
    prs.save(output_file)
    print(f"Generated {output_file}")

if __name__ == "__main__":
    create_ppt()
